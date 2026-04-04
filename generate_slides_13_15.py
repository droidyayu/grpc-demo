from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()

BG = RGBColor(11, 28, 46)
HEADER = RGBColor(39, 63, 91)
ACCENT = RGBColor(0, 187, 221)
TEXT = RGBColor(239, 244, 247)
MUTED = RGBColor(183, 203, 222)


def paint_bg(slide):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_header(slide, title):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(30)
    run.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.45)


# Slide 13
s13 = prs.slides.add_slide(prs.slide_layouts[6])
paint_bg(s13)
add_header(s13, "THE MIGRATION STRATEGY: 5 STEPS, ZERO DOWNTIME")

steps = [
    ("0", "Abstract behind an interface", "ViewModel talks to repository, never to transport."),
    ("1", "Define the .proto", "Mirror REST contract and make proto the source of truth."),
    ("2", "Run gRPC alongside REST", "Add transport in parallel. Do not rip out what works."),
    ("3", "Migrate endpoint by endpoint", "Feature flag, 1% rollout, validate, then expand."),
    ("4", "Deprecate REST gracefully", "Headers + timeline + telemetry before removal."),
]

y = 1.2
for n, title, sub in steps:
    row = s13.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(12.1), Inches(1.0))
    row.fill.solid()
    row.fill.fore_color.rgb = RGBColor(27, 52, 77)
    row.line.color.rgb = RGBColor(16, 38, 58)

    badge = s13.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(0.8), Inches(1.0))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()

    btf = badge.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = n
    br.font.bold = True
    br.font.size = Pt(26)
    br.font.color.rgb = TEXT
    bp.alignment = PP_ALIGN.CENTER

    txt = s13.shapes.add_textbox(Inches(1.55), Inches(y + 0.08), Inches(10.8), Inches(0.85))
    tf = txt.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.bold = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = TEXT
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(13)
    r2.font.color.rgb = MUTED

    y += 1.03

s13.notes_slide.notes_text_frame.text = (
    "Signpost slide: five steps, zero downtime. "
    "We will unpack each step in detail next."
)

# Slide 14
s14 = prs.slides.add_slide(prs.slide_layouts[6])
paint_bg(s14)
add_header(s14, "STEP 0 - Before You Write a Single gRPC Line")

left = s14.shapes.add_shape(1, Inches(0.35), Inches(1.2), Inches(6.1), Inches(5.6))
left.fill.solid()
left.fill.fore_color.rgb = RGBColor(0, 23, 47)
left.line.color.rgb = ACCENT

right = s14.shapes.add_shape(1, Inches(6.75), Inches(1.2), Inches(6.1), Inches(5.6))
right.fill.solid()
right.fill.fore_color.rgb = RGBColor(0, 23, 47)
right.line.color.rgb = RGBColor(235, 158, 86)

interface_code = (
    "interface CatalogRepository {\n"
    "  suspend fun getProduct(productId: String): Result<Product>\n"
    "  fun listProducts(category: String, page: Int): Flow<Product>\n"
    "  fun watchPrice(productId: String): Flow<PriceUpdate>\n"
    "}"
)

swap_code = (
    "@Provides\n"
    "@Singleton\n"
    "fun provideCatalogRepository(\n"
    "  flags: FeatureFlags,\n"
    "  grpcImpl: CatalogRepositoryImpl,\n"
    "  restImpl: RestCatalogRepositoryImpl\n"
    "): CatalogRepository =\n"
    "  if (flags.useGrpcCatalog()) grpcImpl else restImpl"
)

for x, label, color in [
    (0.45, "THE INTERFACE", ACCENT),
    (6.9, "THE SWAP", RGBColor(235, 158, 86)),
]:
    box = s14.shapes.add_textbox(Inches(x), Inches(1.27), Inches(2.2), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = color

left_txt = s14.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.7), Inches(4.9))
ltf = left_txt.text_frame
ltf.clear()
lp = ltf.paragraphs[0]
lr = lp.add_run()
lr.text = interface_code
lr.font.name = "Menlo"
lr.font.size = Pt(18)
lr.font.color.rgb = RGBColor(114, 216, 255)

right_txt = s14.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.9))
rtf = right_txt.text_frame
rtf.clear()
rp = rtf.paragraphs[0]
rr = rp.add_run()
rr.text = swap_code
rr.font.name = "Menlo"
rr.font.size = Pt(17)
rr.font.color.rgb = RGBColor(242, 184, 122)

banner = s14.shapes.add_shape(1, Inches(0), Inches(6.8), prs.slide_width, Inches(0.7))
banner.fill.solid()
banner.fill.fore_color.rgb = ACCENT
banner.line.fill.background()

btf = banner.text_frame
btf.clear()
bp = btf.paragraphs[0]
br = bp.add_run()
br.text = "Your ViewModel does not know what gRPC is. That is the point."
br.font.bold = True
br.font.size = Pt(22)
br.font.color.rgb = RGBColor(7, 24, 39)
bp.alignment = PP_ALIGN.CENTER

s14.notes_slide.notes_text_frame.text = (
    "Step 0 is architecture first: isolate transport behind an interface, "
    "then use DI + feature flag to switch between REST and gRPC safely."
)

# Slide 15
s15 = prs.slides.add_slide(prs.slide_layouts[6])
paint_bg(s15)
add_header(s15, "STEP 1 - DEFINE YOUR CONTRACT")

proto_panel = s15.shapes.add_shape(1, Inches(0.4), Inches(1.15), Inches(7.2), Inches(5.9))
proto_panel.fill.solid()
proto_panel.fill.fore_color.rgb = RGBColor(0, 23, 47)
proto_panel.line.color.rgb = ACCENT

proto_code = (
    "// product_catalog.proto\n"
    "syntax = \"proto3\";\n"
    "package catalog;\n\n"
    "message Product {\n"
    "  string id = 1;\n"
    "  string name = 2;\n"
    "  int64 price_paise = 3;\n"
    "  string image_url = 4;\n"
    "  bool in_stock = 5;\n"
    "}\n\n"
    "service CatalogService {\n"
    "  rpc GetProduct(GetProductRequest) returns (Product);\n"
    "  rpc ListProducts(ListProductsReq) returns (stream Product);\n"
    "}"
)

proto_txt = s15.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(6.8), Inches(5.3))
ptf = proto_txt.text_frame
ptf.clear()
pp = ptf.paragraphs[0]
pr = pp.add_run()
pr.text = proto_code
pr.font.name = "Menlo"
pr.font.size = Pt(18)
pr.font.color.rgb = RGBColor(114, 216, 255)

title = s15.shapes.add_textbox(Inches(7.9), Inches(1.3), Inches(5.0), Inches(0.5))
tp = title.text_frame.paragraphs[0]
tr = tp.add_run()
tr.text = "Run protoc -> you get"
tr.font.bold = True
tr.font.size = Pt(22)
tr.font.color.rgb = TEXT

cards = [
    ("Android (Kotlin)", "ProductOuterClass.kt + gRPC Kotlin stub"),
    ("iOS (Swift)", "product_catalog.pb.swift + product_catalog.grpc.swift"),
    ("Flutter (Dart)", "product_catalog.pb.dart + grpc client types"),
    ("Backend (Ktor)", "Server stubs ready to implement"),
]

y = 1.95
for h, sub in cards:
    card = s15.shapes.add_shape(1, Inches(7.9), Inches(y), Inches(4.8), Inches(1.1))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(29, 57, 85)
    card.line.color.rgb = ACCENT

    txt = s15.shapes.add_textbox(Inches(8.05), Inches(y + 0.14), Inches(4.5), Inches(0.8))
    tf = txt.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = h
    r1.font.bold = True
    r1.font.size = Pt(16)
    r1.font.color.rgb = RGBColor(48, 208, 243)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(13)
    r2.font.color.rgb = MUTED
    y += 1.2

note = s15.shapes.add_shape(1, Inches(7.9), Inches(6.95), Inches(4.8), Inches(0.45))
note.fill.solid()
note.fill.fore_color.rgb = RGBColor(235, 158, 86)
note.line.fill.background()
ntf = note.text_frame
ntf.clear()
np = ntf.paragraphs[0]
nr = np.add_run()
nr.text = "Versioning: never renumber tags, never reuse removed fields"
nr.font.bold = True
nr.font.size = Pt(11)
nr.font.color.rgb = RGBColor(21, 26, 31)
np.alignment = PP_ALIGN.CENTER

s15.notes_slide.notes_text_frame.text = (
    "Contract first. Proto is one source of truth across platforms. "
    "Field numbers are wire IDs, so keep compatibility rules strict."
)

out_path = "grpc_migration_slides_13_15.pptx"
prs.save(out_path)
print(out_path)
